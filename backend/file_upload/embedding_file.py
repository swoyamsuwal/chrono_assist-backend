# ===============================================================
#  file_upload/embedding_file.py
#  The full RAG embedding pipeline for a single Document
#
#  FLOW OVERVIEW:
#  Step 1 → Read file from MinIO via Django FileField
#  Step 2 → Extract raw text (PDF / DOCX / PPTX / TXT)
#  Step 3 → Split text into overlapping chunks
#  Step 4 → Generate vector embeddings per chunk via Ollama
#  Step 5 → Save all chunks + vectors to DocumentChunk (pgvector)
#  Step 6 → Mark Document.is_embedded = True
# ===============================================================


# ---------------- Step 0: Imports & Config ----------------
import os
import tempfile

# File parsers for each supported format
from pypdf import PdfReader                # Extracts text page-by-page from PDF
from docx import Document as DocxDocument  # Reads DOCX paragraph objects
from pptx import Presentation              # Reads PPTX slide shapes

# LangChain splitter — handles smart splitting that respects word/sentence boundaries
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Ollama Python client — used to call the local embedding model
from ollama import Client

# Django models
from .models import Document, DocumentChunk

# Ollama host — read from environment so it works in Docker or local dev
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://localhost:11434")
# The embedding model must be pulled in Ollama before this runs: `ollama pull all-minilm:l6-v2`
EMBEDDING_MODEL_NAME = "all-minilm:l6-v2"


# ================================================================
#  Function 1: extract_text_from_fileobj
#  Entry point for reading a Django FileField
#
#  Problem: Django FileField (backed by MinIO) gives us a file-like object,
#  but PDF/DOCX/PPTX parsers require a real path on disk.
#  Solution: Stream the file into a temp file, parse the temp file, then delete it.
# ================================================================
def extract_text_from_fileobj(django_file, mime_type: str) -> str:

    # ---------------- Step 1a: Determine File Extension ----------------
    # Used as the suffix for the temp file so parsers can identify the format
    suffix = "." + (django_file.name.split(".")[-1] if "." in django_file.name else "")

    # ---------------- Step 1b: Write to Temp File ----------------
    # django_file.chunks() streams the file in memory-safe blocks (avoids loading all at once)
    # delete=False → we manage cleanup manually (in the finally block below)
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        for chunk in django_file.chunks():
            tmp.write(chunk)
        temp_path = tmp.name  # e.g., /tmp/tmpXYZ.pdf

    # ---------------- Step 1c: Parse + Cleanup ----------------
    try:
        return extract_text_from_path(temp_path, mime_type)
    finally:
        # Always delete the temp file even if parsing fails
        try:
            os.remove(temp_path)
        except OSError:
            pass


# ================================================================
#  Function 2: extract_text_from_path
#  Reads a local file path and extracts all readable text
#  Dispatches to the correct parser based on mime_type or extension
# ================================================================
def extract_text_from_path(file_path: str, mime_type: str) -> str:
    mime_type = (mime_type or "").lower()
    ext = os.path.splitext(file_path)[1].lower()  # e.g., ".pdf", ".docx"

    # ---------------- Step 2a: PDF Parsing ----------------
    # PdfReader extracts text page by page — joins all pages with newline
    if "pdf" in mime_type or ext == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    # ---------------- Step 2b: DOCX Parsing ----------------
    # DocxDocument reads paragraphs — joins them with newline
    if "word" in mime_type or ext in (".docx",):
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    # ---------------- Step 2c: PPTX Parsing ----------------
    # Iterates every slide and every shape on the slide that has text
    if "ppt" in mime_type or ext in (".pptx",):
        pres = Presentation(file_path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    # ---------------- Step 2d: TXT / Fallback ----------------
    # Everything else is treated as plain text
    # errors="ignore" skips unreadable bytes instead of crashing
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ================================================================
#  Function 3: chunk_text
#  Splits a long text into overlapping segments
#
#  Why overlap? So that a sentence spanning a chunk boundary
#  still appears in at least one chunk's context window
#  chunk_size=1000 → ~750 words per chunk (good for LLaMA context)
#  chunk_overlap=200 → 200 chars of the previous chunk repeated at the start of the next
# ================================================================
def chunk_text(text: str):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,    # Max characters per chunk
        chunk_overlap=200,  # Characters shared between adjacent chunks
        length_function=len,
    )
    return splitter.split_text(text)  # Returns list[str]


# ================================================================
#  Function 4: embed_chunks
#  Sends each chunk to Ollama and collects the embedding vectors
#
#  all-minilm:l6-v2 → produces 384-dimensional float vectors
#  Each vector represents the semantic meaning of that chunk in high-dimensional space
#  CosineDistance can then find the most semantically similar chunks for any query
# ================================================================
def embed_chunks(chunks):
    # Client() reads OLLAMA_HOST from environment internally
    client = Client()

    embeddings = []
    for chunk in chunks:
        # client.embeddings() returns {"embedding": [float, float, ...]}
        resp = client.embeddings(model=EMBEDDING_MODEL_NAME, prompt=chunk)
        embeddings.append(resp["embedding"])

    return embeddings  # list[list[float]] — one vector per chunk


# ================================================================
#  Function 5: create_embeddings_for_document  (MAIN ENTRY POINT)
#  Orchestrates the full pipeline for one Document
#  Called by the embed_file view when the user clicks "embed"
#
#  Full pipeline:
#   1. Read file from MinIO (via Django FileField)
#   2. Extract raw text (format-aware)
#   3. Split text into chunks
#   4. Generate one embedding vector per chunk via Ollama
#   5. Delete old chunks (re-embed support) and bulk-insert new ones
#   6. Mark Document.is_embedded = True
# ================================================================
def create_embeddings_for_document(doc: Document):

    # ---------------- Step 5a: Extract Text from File ----------------
    # doc.file is the Django FileField — stored in MinIO, accessed via streaming
    django_file = doc.file
    text = extract_text_from_fileobj(django_file, doc.mime_type or "")

    # If the file had no parseable text (e.g., scanned image PDF), stop early
    if not text.strip():
        return 0  # 0 chunks created

    # ---------------- Step 5b: Chunk the Text ----------------
    chunks = chunk_text(text)

    # ---------------- Step 5c: Generate Embeddings ----------------
    vectors = embed_chunks(chunks)

    # ---------------- Step 5d: Save Chunks to DB ----------------
    # Delete existing chunks first — supports re-embedding if the file changes
    DocumentChunk.objects.filter(document=doc).delete()

    objs = []
    for idx, (chunk_text_value, vec) in enumerate(zip(chunks, vectors)):
        objs.append(
            DocumentChunk(
                document=doc,
                chunk_index=idx,           # Position of this chunk in the document
                text=chunk_text_value,     # The raw text shown as context to the LLM
                embedding=vec,             # The pgvector float[] used for similarity search
                embedding_model=EMBEDDING_MODEL_NAME,
            )
        )

    # bulk_create inserts all rows in one SQL statement — much faster than individual saves
    DocumentChunk.objects.bulk_create(objs)

    # ---------------- Step 5e: Mark Document as Embedded ----------------
    # is_embedded = True unlocks rag_chat and doc_chat for this document
    doc.is_embedded = True
    doc.save(update_fields=["is_embedded", "updated_at"])

    return len(objs)  # Returned to the API response as "chunks_created"